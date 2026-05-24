from __future__ import annotations  # 启用延迟类型注解，避免运行时前向引用问题

from pathlib import Path  # 导入 Path 以统一处理文件路径

from app.documents.loaders.base import LoadedPage  # 导入文档解析或切片相关组件
from app.text_utils import clean_text  # 导入文本清洗工具以移除非法字符


class PptLoader:  # 按 Slide 提取 PPT/PPTX 文本
    def load(self, path: str | Path) -> list[LoadedPage]:  # 定义 load 函数或方法
        try:  # 尝试执行可能失败的操作
            from pptx import Presentation  # 导入当前模块需要的依赖
        except ImportError as exc:  # 捕获依赖未安装的情况
            raise RuntimeError("python-pptx is required to load PPT/PPTX files. Run: pip install -r requirements.txt") from exc  # 将缺失依赖转换成可读的运行时错误

        presentation = Presentation(str(path))  # 计算并保存 presentation
        pages: list[LoadedPage] = []  # 保存当前流程中的 pages 数据
        for index, slide in enumerate(presentation.slides, start=1):  # 遍历当前集合中的每个元素
            texts: list[str] = []  # 计算并保存 texts
            for shape in slide.shapes:  # 遍历当前集合中的每个元素
                if hasattr(shape, "text") and shape.text:  # 检查条件：hasattr(shape, "text") and shape.text
                    texts.append(shape.text)  # 将当前结果追加到列表
            text = clean_text("\n".join(texts))  # 保存切片原文内容
            if text.strip():  # 检查条件：text.strip()
                pages.append(LoadedPage(text=text, slide_number=index))  # 将当前结果追加到列表
        return pages  # 返回 pages
